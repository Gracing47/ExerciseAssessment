from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Team Availability Analysis"
    subtitle.text = "August 2022\nAnalysis of staffing requirements and recommendations"

def create_requirements_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Service Line Requirements"
    tf = body.text_frame
    tf.text = "Minimum Staffing Requirements:"
    
    p = tf.add_paragraph()
    p.text = "• Service Line 1: 2.5 agents"
    p = tf.add_paragraph()
    p.text = "• Service Line 2: 4.0 agents"
    p = tf.add_paragraph()
    p.text = "• Service Line 3: 18.0 agents"
    p = tf.add_paragraph()
    p.text = "\nTotal Required: 24.5 agents"

def create_statistics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Statistics"
    tf = body.text_frame
    
    stats = [
        "Total Agents: 30",
        "Working Days: 23 weekdays (Monday-Friday only)",
        "Average Daily Available Agents: 25.08",
        "Maximum Available: 29.00 agents",
        "Minimum Available: 21.00 agents",
        "Days Meeting Requirements: 16 days (≥24.5 agents)",
        "Days Below Requirements: 7 days (<24.5 agents)"
    ]
    
    for i, stat in enumerate(stats):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {stat}"

def create_analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Staffing Analysis"
    tf = body.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Coverage Rate:"
    p = tf.add_paragraph()
    p.text = "• 69.6% of weekdays met minimum requirement"
    p = tf.add_paragraph()
    p.text = "• 30.4% of weekdays fell below requirement"
    
    p = tf.add_paragraph()
    p.text = "\nStaffing Buffer:"
    p = tf.add_paragraph()
    p.text = "• Average surplus: 0.58 agents"
    p = tf.add_paragraph()
    p.text = "• Maximum surplus: 4.5 agents"
    p = tf.add_paragraph()
    p.text = "• Maximum deficit: 3.5 agents"

def create_critical_issues_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Critical Issues"
    tf = body.text_frame
    
    issues = [
        ("Staffing Gaps:", [
            "Highest gaps in third week of August",
            "Friday staffing levels show shortages",
            "7 days (30.4%) below requirements"
        ]),
        ("Resource Management:", [
            "Uneven distribution of leave days",
            "Limited buffer for unexpected absences",
            "Peak vacation period impacts"
        ])
    ]
    
    for i, (header, points) in enumerate(issues):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = header
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"

def create_recommendations_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Recommendations"
    tf = body.text_frame
    
    recommendations = [
        ("Leave Management:", [
            "Implement structured approval process",
            "Stagger vacation schedules",
            "Set concurrent leave limits"
        ]),
        ("Staffing Adjustments:", [
            "Add 1-2 additional agents",
            "Develop flexible staffing model",
            "Cross-train across service lines"
        ]),
        ("Monitoring System:", [
            "Implement daily tracking",
            "Set up shortage alerts",
            "Regular pattern analysis"
        ])
    ]
    
    for i, (header, points) in enumerate(recommendations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = header
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"

def create_action_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Action Plan"
    tf = body.text_frame
    
    actions = [
        ("Immediate (1-2 weeks):", [
            "Review and optimize staff distribution",
            "Implement leave request coordination"
        ]),
        ("Short-term (1-3 months):", [
            "Begin recruitment for additional staff",
            "Develop cross-training program"
        ]),
        ("Long-term (3-6 months):", [
            "Implement workforce management system",
            "Establish monitoring process"
        ])
    ]
    
    for i, (header, points) in enumerate(actions):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = header
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"

def create_visualization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    
    title.text = "Availability Analysis"
    
    # Add daily availability chart
    img_path = os.path.join('..', 'output', 'daily_availability.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

def create_presentation():
    prs = Presentation()
    
    # Create slides
    create_title_slide(prs)
    create_requirements_slide(prs)
    create_statistics_slide(prs)
    create_analysis_slide(prs)
    create_critical_issues_slide(prs)
    create_recommendations_slide(prs)
    create_action_plan_slide(prs)
    create_visualization_slide(prs)
    
    # Save presentation
    output_dir = os.path.join('..', 'output', 'Presentation')
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs.save(os.path.join(output_dir, 'team_availability_analysis.pptx'))

if __name__ == '__main__':
    create_presentation()
