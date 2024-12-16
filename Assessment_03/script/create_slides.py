from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Team Availability Analysis"
    subtitle.text = "August 2022\nStaffing Analysis & Recommendations"

def create_requirements_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Service Line Requirements"
    tf = body.text_frame
    tf.text = "Minimum Staffing Requirements:"
    
    p = tf.add_paragraph()
    p.text = "• Service Line 1: 2.5 agents"
    p = tf.add_paragraph()
    p.text = "• Service Line 2: 4.0 agents"
    p = tf.add_paragraph()
    p.text = "• Service Line 3: 18.0 agents"
    p = tf.add_paragraph()
    p.text = "\nTotal Required: 24.5 agents"

def create_statistics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Statistics"
    tf = body.text_frame
    
    stats = [
        "Total Agents: 30",
        "Working Days: 23 weekdays (Monday-Friday only)",
        "Average Daily Available Agents: 25.08",
        "Maximum Available: 29.00 agents",
        "Minimum Available: 21.00 agents",
        "Days Meeting Requirements: 16 days (≥24.5 agents)",
        "Days Below Requirements: 7 days (<24.5 agents)"
    ]
    
    for i, stat in enumerate(stats):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {stat}"

def create_sick_leave_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Sick Leave Analysis"
    tf = body.text_frame
    
    stats = [
        "Total sick leave instances: 28",
        "Maximum sick leaves in one day: 4",
        "Average sick leaves per day: 1.22",
        "Days with sick leave: 14 (60.9%)"
    ]
    
    for i, stat in enumerate(stats):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"• {stat}"
    
    # Add sick leave visualization
    img_path = os.path.join('..', 'output', 'sick_leave_analysis.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(3), width=Inches(8))

def create_recommendations_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Recommendations"
    tf = body.text_frame
    
    recommendations = [
        ("Leave Management:", [
            "Structured approval process",
            "Stagger vacation schedules",
            "Monitor sick leave patterns"
        ]),
        ("Staffing Adjustments:", [
            "Add 1-2 additional agents",
            "Cross-train across service lines",
            "Implement wellness programs"
        ]),
        ("Monitoring:", [
            "Daily tracking system",
            "Early warning alerts",
            "Regular pattern analysis"
        ])
    ]
    
    for i, (header, points) in enumerate(recommendations):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = header
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"

def create_action_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Action Plan"
    tf = body.text_frame
    
    actions = [
        ("Now (1-2 weeks):", [
            "Optimize staff distribution",
            "Review sick leave patterns"
        ]),
        ("Soon (1-3 months):", [
            "Start recruitment",
            "Launch wellness program"
        ]),
        ("Later (3-6 months):", [
            "Implement new systems",
            "Review effectiveness"
        ])
    ]
    
    for i, (header, points) in enumerate(actions):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = header
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"

def create_visualization_slides(prs):
    # Availability Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Daily Availability"
    
    img_path = os.path.join('..', 'output', 'daily_availability.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))
    
    # Staffing Gaps Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Staffing Gaps"
    
    img_path = os.path.join('..', 'output', 'staffing_gaps.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

def create_presentation():
    prs = Presentation()
    
    # Create slides
    create_title_slide(prs)
    create_requirements_slide(prs)
    create_statistics_slide(prs)
    create_sick_leave_slide(prs)
    create_recommendations_slide(prs)
    create_action_plan_slide(prs)
    create_visualization_slides(prs)
    
    # Save presentation with shorter name
    output_dir = os.path.join('..', 'output', 'Presentation')
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs.save(os.path.join(output_dir, 'team_analysis.pptx'))

if __name__ == '__main__':
    create_presentation()
