from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Formatierung
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    
    return slide

def create_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Titel
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    
    for line in content.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
    
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 Format
    prs.slide_height = Inches(9)
    
    # Titelfolie
    create_title_slide(prs, 
                      "Loan Application Process Analysis 🏦",
                      "Bottlenecks and Solutions Study")
    
    # Application Steps
    steps_content = """
Key Process Steps:

1. Initial Application
   • Basic personal information
   • Loan amount and purpose
   • Initial eligibility check

2. Documentation
   • Income verification
   • Employment history
   • Bank statements
   • Tax returns

3. Credit Assessment
   • Credit score check
   • Debt-to-income ratio
   • Payment history analysis
   • Risk assessment
"""
    create_content_slide(prs, "Application Steps 📋", steps_content)
    
    # More Steps
    more_steps_content = """
4. Property Evaluation
   • Property documentation
   • Appraisal process
   • Title search
   • Insurance requirements

5. Final Review
   • Underwriting process
   • Terms finalization
   • Approval/rejection decision
   • Offer presentation
"""
    create_content_slide(prs, "Application Steps (continued) 📋", more_steps_content)
    
    # Major Bottlenecks
    bottlenecks_content = """
Critical Pain Points:

1. Documentation Phase (40% Drop-off)
   • Excessive document requirements
   • Complex submission process
   • Multiple format requirements
   • Time-consuming gathering process

2. Credit Assessment (25% Drop-off)
   • Long waiting periods
   • Lack of transparency
   • Unclear requirements
   • Poor communication

3. Property Evaluation (20% Drop-off)
   • Scheduling delays
   • Coordination issues
   • Multiple visits required
   • Price negotiations
"""
    create_content_slide(prs, "Major Bottlenecks 🚧", bottlenecks_content)
    
    # Impact Analysis
    impact_content = """
Drop-off Statistics:

Initial Application:
• 100% start the process
• 85% complete basic information
• 60% reach documentation phase

Documentation Phase:
• 60% start documentation
• Only 36% complete all documents
• 20% drop due to complexity

Final Stages:
• 35% reach final review
• 25% receive approval
• 20% complete the process
"""
    create_content_slide(prs, "Impact Analysis 📊", impact_content)
    
    # Root Causes
    causes_content = """
Key Issues Identified:

Documentation Challenges:
• Multiple document formats
• Repetitive information requests
• Unclear requirements
• Manual verification processes

Process Issues:
• Long waiting times
• Poor communication
• Complex requirements
• Limited digital integration

Customer Experience:
• Lack of transparency
• Multiple touchpoints
• Inconsistent information
• Limited self-service options
"""
    create_content_slide(prs, "Root Causes 🔍", causes_content)
    
    # Proposed Solutions
    solutions_content = """
Digital Transformation:

1. Smart Documentation
   • Digital document upload
   • Auto-format conversion
   • Real-time validation
   • Document pre-filling

2. Process Automation
   • Automated eligibility checks
   • Real-time status updates
   • Integrated credit checks
   • Smart underwriting

3. Customer Experience
   • Mobile-first approach
   • Progress tracking
   • Interactive guides
   • Chat support
"""
    create_content_slide(prs, "Proposed Solutions ✨", solutions_content)
    
    # Implementation Plan
    implementation_content = """
Phased Approach:

Phase 1: Digital Foundation
• Mobile application platform
• Document upload system
• Basic automation

Phase 2: Process Enhancement
• Automated assessments
• Integration with credit bureaus
• Real-time status tracking

Phase 3: Advanced Features
• AI-powered pre-approval
• Predictive analytics
• Full process automation
"""
    create_content_slide(prs, "Implementation Plan 📈", implementation_content)
    
    # Expected Improvements
    improvements_content = """
Key Metrics:

Completion Rates:
• Documentation: +40%
• Credit Assessment: +30%
• Overall Process: +35%

Time Reduction:
• Application: -50%
• Processing: -60%
• Total Time: -55%

Customer Satisfaction:
• Experience: +45%
• Recommendation: +50%
• Retention: +40%
"""
    create_content_slide(prs, "Expected Improvements 📊", improvements_content)
    
    # Conclusion
    conclusion_content = """
Key Takeaways:

1. Major bottlenecks identified in:
   • Documentation collection
   • Credit assessment
   • Property evaluation

2. Solutions focused on:
   • Digital transformation
   • Process automation
   • Customer experience

3. Expected outcomes:
   • Higher completion rates
   • Faster processing
   • Improved satisfaction
"""
    create_content_slide(prs, "Conclusion 🎯", conclusion_content)
    
    # Speichern
    output_dir = "output/Presentation"
    os.makedirs(output_dir, exist_ok=True)
    prs.save(f"{output_dir}/loan_application_analysis.pptx")

if __name__ == "__main__":
    create_presentation()
