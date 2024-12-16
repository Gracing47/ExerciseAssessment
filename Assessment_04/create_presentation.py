from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Formatierung
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    
    return slide

def create_content_slide(prs, title, content, has_image=False, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Titel
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    
    # Content
    if has_image and image_path and os.path.exists(image_path):
        # Bild links, Text rechts
        img_left = Inches(1)
        img_top = Inches(2)
        img_width = Inches(4)
        img_height = Inches(4)
        
        slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
        
        # Text rechts vom Bild
        text_left = Inches(6)
        text_top = Inches(2)
        text_width = Inches(4)
        text_height = Inches(4)
        
        textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        text_frame = textbox.text_frame
        
        for line in content.split('\n'):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
    else:
        # Nur Text
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        
        for line in content.split('\n'):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
    
    return slide

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 Format
    prs.slide_height = Inches(9)
    
    # Titelfolie
    create_title_slide(prs, 
                      "Call Efficiency Analysis 📞",
                      "Performance vs. Call Duration Study")
    
    # Overview
    overview_content = """
Analysis Scope:
• Focus: Call Duration vs. Efficiency
• Number of Agents: 4
• Key Metrics:
  - Call Duration
  - Calls per Day
  - Overall Efficiency
"""
    create_content_slide(prs, "Overview 📊", overview_content)
    
    # Call Duration Analysis
    duration_content = """
Average Call Times:
• Agent 1: 237 seconds
• Agent 2: 194 seconds
• Agent 3: 162 seconds
• Agent 4: 201 seconds
• Team Average: 198.5 seconds
"""
    create_content_slide(prs, "Call Duration Analysis ⏱️", duration_content, True,
                        "output/call_duration_by_agent.png")
    
    # Daily Call Volume
    volume_content = """
Calls per Day:
• Agent 1: 60 calls/day
• Agent 2: 80 calls/day
• Agent 3: 110 calls/day
• Agent 4: 70 calls/day
• Team Average: 80 calls/day
"""
    create_content_slide(prs, "Daily Call Volume 📈", volume_content, True,
                        "output/calls_per_day_by_agent.png")
    
    # Efficiency Comparison
    efficiency_content = """
Performance Analysis:

Fast Calls (Agent 3):
• Shortest duration: 162 seconds
• Highest volume: 110 calls/day
• Maximum efficiency in terms of quantity

Longer Calls (Agent 1):
• Longest duration: 237 seconds
• Lower volume: 60 calls/day
• Focus on quality over quantity
"""
    create_content_slide(prs, "Efficiency Comparison 📊", efficiency_content)
    
    # Key Findings
    findings_content = """
Data Insights:
1. Clear trade-off between call duration and volume
2. Agent 3 processes 83% more calls than Agent 1
3. Call duration varies by up to 75 seconds
4. Significant impact on daily productivity
"""
    create_content_slide(prs, "Key Findings 🔍", findings_content)
    
    # Efficiency Factors
    factors_content = """
Considerations:

Shorter Calls:
• Higher call volume
• More customers served
• Reduced wait times

Longer Calls:
• Detailed customer service
• Potential for better resolution
• Higher customer satisfaction
"""
    create_content_slide(prs, "Efficiency Factors ⚖️", factors_content)
    
    # Recommendations
    recommendations_content = """
Balanced Approach:
1. Set optimal duration targets
2. Consider call complexity
3. Balance quantity with quality
4. Implement targeted training
"""
    create_content_slide(prs, "Recommendations ✨", recommendations_content)
    
    # Conclusion
    conclusion_content = """
Optimal Strategy:
• Agent 3's approach shows highest efficiency
• Focus on reducing call duration while maintaining quality
• Implement best practices from efficient agents
• Regular performance monitoring
"""
    create_content_slide(prs, "Conclusion 🎯", conclusion_content)
    
    # Speichern
    output_dir = "output/Presentation"
    os.makedirs(output_dir, exist_ok=True)
    prs.save(f"{output_dir}/call_efficiency_analysis.pptx")

if __name__ == "__main__":
    create_presentation()
